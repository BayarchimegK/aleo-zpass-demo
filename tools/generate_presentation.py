from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "ALEO zPass — Age Verification Demo"
subtitle.text = f"Presentation generated on {datetime.now().strftime('%Y-%m-%d')}"

# Overview slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
sh = slide.shapes.title
sh.text = "Architecture Overview"
body = slide.shapes.placeholders[1].text_frame
body.text = "Three-tier architecture"

p = body.add_paragraph()
p.level = 1
p.text = "• Frontend: Next.js (port 3000) — holder & verifier UIs"

p = body.add_paragraph()
p.level = 1
p.text = "• Auth Backend: Express + Prisma (port 4000) — nonce, proof, JWTs"

p = body.add_paragraph()
p.level = 1
p.text = "• Content Backend: Express (port 4001) — serves gated content, checks revocation"

# ZKP Workflow slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "ZKP Age Verification Workflow"
body = slide.shapes.placeholders[1].text_frame
body.text = "High-level flow (nonce, proof, JWT)"

steps = [
    "1. Browser requests nonce (GET /auth/nonce). Server stores 1-use nonce (5min).",
    "2. Client obtains credential metadata by email (GET /issuer/credential-by-email).",
    "3. Client POSTs email+nonce to /auth/verify — backend validates nonce and revocation.",
    "4. Backend runs Leo locally or via SSH: `leo run prove_age_over_18 ${age}u8`.",
    "5. On success, backend issues short-lived JWT (15m) with claim `isAdult`.",
    "6. Frontend requests content from Content Backend (JWT attached). Content backend re-checks revocation.",
]

for s in steps:
    p = body.add_paragraph()
    p.level = 1
    p.text = s

# Security properties slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Security Properties"
body = slide.shapes.placeholders[1].text_frame
body.text = "Key mitigations"

props = [
    "One-time nonce (5-minute TTL) to prevent replay attacks.",
    "Revocation checked before proof generation and on every content request.",
    "Short-lived JWT (15 minutes) to limit token exposure.",
    "Rate limiting on /auth/verify (5/min per IP).",
    "Leo proves boolean `age >= 18` — exact age never leaves the prover."
]
for s in props:
    p = body.add_paragraph()
    p.level = 1
    p.text = s

# Data models slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Database Models (Prisma)"
body = slide.shapes.placeholders[1].text_frame
body.text = "Main tables"

models = [
    "Credential: holderEmail, age, country, isRevoked, proofTxId, issuerSignature",
    "Nonce: id (uuid), value, used (bool), expiresAt",
    "User: issuer accounts (admin/issuer)"
]
for m in models:
    p = body.add_paragraph()
    p.level = 1
    p.text = m

# Endpoints slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Key Endpoints"
body = slide.shapes.placeholders[1].text_frame
body.text = "Auth Backend"

endpoints = [
    "GET /auth/nonce — returns one-time nonce",
    "POST /auth/verify — runs proof, returns SSE stream and short JWT",
    "GET /issuer/credential-by-email?email= — returns credential metadata (no age)",
    "GET /auth/check-revocation/:credentialId — revocation status"
]
for e in endpoints:
    p = body.add_paragraph()
    p.level = 1
    p.text = e

# Commands slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Developer Commands"
body = slide.shapes.placeholders[1].text_frame
body.text = "Useful commands"

cmds = [
    "Run backend (from backend/): npx ts-node --transpile-only src/index.ts",
    "Run content backend (from backend-content/): npm run dev",
    "Run frontend (from frontend/): npm run dev",
    "Leo proof (local or SSH): leo run prove_age_over_18 25u8"
]
for c in cmds:
    p = body.add_paragraph()
    p.level = 1
    p.text = c

# Content gating slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Content Gating"
body = slide.shapes.placeholders[1].text_frame
body.text = "JWT claim `isAdult` determines audience"

p = body.add_paragraph()
p.level = 1
p.text = "isAdult=true → Adult tiles: Premium, Marketplace, Finance, Gaming"

p = body.add_paragraph()
p.level = 1
p.text = "isAdult=false → Minor tiles: Education, Youth Games, Story World, Creative Studio"

# Next steps / Review slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Next Steps"
body = slide.shapes.placeholders[1].text_frame
body.text = "Review and adjust for audience"

p = body.add_paragraph()
p.level = 1
p.text = "— Add architecture diagram image or sequence diagram"

p = body.add_paragraph()
p.level = 1
p.text = "— Tailor language for management vs technical audiences"

# Save
out_path = "docs/aleo-zpass-demo-presentation.pptx"
prs.save(out_path)
print(f"Saved presentation to {out_path}")
